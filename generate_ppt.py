"""
PharmaTrust Hackathon Presentation Generator
Run: pip install python-pptx && python generate_ppt.py
Output: PharmaTrust_Hackathon.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color Palette ──────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT     = RGBColor(0x00, 0xD4, 0xFF)   # Cyan
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)   # Purple
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xC4, 0xDE)
GREEN      = RGBColor(0x00, 0xE5, 0x96)
RED        = RGBColor(0xFF, 0x4D, 0x4D)
CARD_BG    = RGBColor(0x16, 0x2A, 0x40)   # Slightly lighter navy for cards

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, alpha=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def accent_bar(slide, t=0.08, h=0.06):
    """Thin gradient-like accent bar at top."""
    box(slide, 0, t, 6.67, h, ACCENT)
    box(slide, 6.67, t, 6.66, h, ACCENT2)

def slide_number(slide, n, total=18):
    txt(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.35,
        size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def bullet_block(slide, items, l, t, w, h, size=15, gap=0.38):
    """Render a list of bullet strings."""
    for i, item in enumerate(items):
        txt(slide, f"  {item}", l, t + i * gap, w, gap + 0.05,
            size=size, color=WHITE)

def card(slide, l, t, w, h, title, body_lines, title_color=ACCENT):
    """Rounded-ish card with title + bullets."""
    box(slide, l, t, w, h, CARD_BG)
    # left accent strip
    box(slide, l, t, 0.04, h, title_color)
    txt(slide, title, l + 0.12, t + 0.08, w - 0.2, 0.35,
        size=14, bold=True, color=title_color)
    for i, line in enumerate(body_lines):
        txt(slide, f"• {line}", l + 0.12, t + 0.48 + i * 0.32,
            w - 0.2, 0.32, size=12, color=LIGHT_GRAY)


def flow_box(slide, l, t, w, h, label, color=ACCENT, text_color=DARK_BG, size=11):
    """Filled rectangle with centered text — used for flowchart nodes."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shape

def arrow_down(slide, l, t, color=LIGHT_GRAY):
    """Small downward arrow for flowcharts."""
    txt(slide, "▼", l, t, 0.5, 0.3, size=13, color=color, align=PP_ALIGN.CENTER)

def arrow_right(slide, l, t, color=LIGHT_GRAY):
    """Small rightward arrow for flowcharts."""
    txt(slide, "▶", l, t, 0.4, 0.3, size=13, color=color, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
# big gradient strip
box(s, 0, 0, 13.33, 0.18, ACCENT)
box(s, 6.67, 0, 6.66, 0.18, ACCENT2)

txt(s, "PharmaTrust AI", 1, 1.4, 11, 1.4,
    size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Pharmaceutical Supply Chain Security Platform",
    1, 2.9, 11, 0.7, size=24, color=ACCENT, align=PP_ALIGN.CENTER)

# divider
box(s, 3.5, 3.75, 6.33, 0.04, ACCENT2)

txt(s, "Blockchain  ·  AI Anomaly Detection  ·  QR Anti-Cloning  ·  Cloud-Native",
    1, 3.95, 11, 0.5, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "Hackathon 2026", 1, 5.5, 11, 0.5,
    size=18, color=ACCENT2, align=PP_ALIGN.CENTER)
slide_number(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "The Problem", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

stats = [
    ("$200B+",  "Annual loss from counterfeit drugs globally"),
    ("1 in 10",  "Medicines in developing countries are substandard"),
    ("~1M",      "Deaths per year linked to fake pharmaceuticals"),
    ("Manual",   "Current supply chain tracking is paper-based & fragile"),
]
for i, (stat, desc) in enumerate(stats):
    col = i % 2
    row = i // 2
    lx = 0.5 + col * 6.5
    ty = 1.4 + row * 2.3
    box(s, lx, ty, 6.1, 2.0, CARD_BG)
    box(s, lx, ty, 0.06, 2.0, RED)
    txt(s, stat, lx + 0.2, ty + 0.15, 5.7, 0.9,
        size=38, bold=True, color=RED)
    txt(s, desc, lx + 0.2, ty + 1.05, 5.7, 0.8,
        size=14, color=LIGHT_GRAY)

slide_number(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Our Solution", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

pillars = [
    (ACCENT,  "🔗 Blockchain",    ["Immutable batch records", "Merkle root verification", "99.99% gas cost reduction"]),
    (ACCENT2, "🤖 AI Sentinel",   ["Impossible travel detection", "Auto kill-switch", "Real-time risk scoring"]),
    (GREEN,   "📱 QR Anti-Clone", ["Protobuf compact payload", "ECDSA digital signature", "Offline verification"]),
    (RGBColor(0xFF,0xA5,0x00), "☁️ Cloud-Native", ["AWS S3 lab reports", "PostgreSQL production DB", "RabbitMQ async jobs"]),
]
for i, (color, title, bullets) in enumerate(pillars):
    lx = 0.4 + i * 3.15
    box(s, lx, 1.3, 2.9, 5.5, CARD_BG)
    box(s, lx, 1.3, 2.9, 0.06, color)
    txt(s, title, lx + 0.1, 1.45, 2.7, 0.5,
        size=15, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", lx + 0.1, 2.05 + j * 0.55, 2.7, 0.5,
            size=12, color=LIGHT_GRAY)

slide_number(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "System Architecture", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

layers = [
    (ACCENT,  "CLIENT LAYER",   "Manufacturer · Distributor · Pharmacist · Regulator"),
    (ACCENT2, "API GATEWAY",    "Spring Boot · JWT Auth Filter · Rate Limiting"),
    (GREEN,   "SERVICE LAYER",  "Auth · Batch · QR · Crypto · AI Sentinel · Storage · Blockchain"),
    (RGBColor(0xFF,0xA5,0x00), "DATA LAYER", "PostgreSQL · AWS S3 · Blockchain Network"),
]
for i, (color, layer, desc) in enumerate(layers):
    ty = 1.4 + i * 1.35
    box(s, 0.5, ty, 12.33, 1.15, CARD_BG)
    box(s, 0.5, ty, 0.06, 1.15, color)
    txt(s, layer, 0.75, ty + 0.1, 3.5, 0.4,
        size=13, bold=True, color=color)
    txt(s, desc, 0.75, ty + 0.55, 11.5, 0.45,
        size=13, color=LIGHT_GRAY)
    # arrow down (except last)
    if i < 3:
        txt(s, "▼", 6.3, ty + 1.15, 0.7, 0.25,
            size=14, color=ACCENT, align=PP_ALIGN.CENTER)

slide_number(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Flowchart: System Architecture Visual  (matches design.md §2.1)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "System Architecture — Flow", 0.5, 0.3, 12, 0.7,
    size=32, bold=True, color=WHITE)

# ── Client Layer ─────────────────────────────────────────────────────────────
# MFG, DIST, PHARM, REG  (design.md: MFG→APIGW, DIST→APIGW, PHARM→APIGW, REG→APIGW)
clients = [
    (ACCENT,                   "🏭 Manufacturer\nDashboard"),
    (ACCENT2,                  "🚚 Distributor\nApp"),
    (GREEN,                    "💊 Pharmacist\nScanner"),
    (RGBColor(0xFF,0xA5,0x00), "🏛️ Regulator\nPortal"),
]
for i, (color, label) in enumerate(clients):
    lx = 0.5 + i * 3.1
    flow_box(s, lx, 1.05, 2.7, 0.65, label, color, DARK_BG, size=11)

# arrows down → API Gateway
for i in range(4):
    arrow_down(s, 1.65 + i * 3.1, 1.73)

# ── API Gateway Layer ─────────────────────────────────────────────────────────
# APIGW + JWT  (design.md: APIGW→JWT, JWT→AUTH/BATCH/QR)
flow_box(s, 0.5, 2.0, 5.9, 0.52,
         "⚡ Spring Boot API Gateway", CARD_BG, ACCENT, size=12)
flow_box(s, 6.9, 2.0, 5.9, 0.52,
         "🔐 JWT Authentication Filter", CARD_BG, ACCENT2, size=12)
arrow_right(s, 6.45, 2.1)

# arrows down from JWT → Service Layer
for lx in [1.25, 4.35, 9.75]:
    arrow_down(s, lx, 2.55)

# ── Service Layer ─────────────────────────────────────────────────────────────
# AUTH, BATCH, QR, CRYPTO, AI, STORAGE, BLOCKCHAIN  (design.md service layer)
services = [
    (ACCENT,                   0.5,  "AUTH\nService"),
    (ACCENT2,                  3.6,  "BATCH\nService"),
    (GREEN,                    6.7,  "QR Code\nService"),
    (RGBColor(0xFF,0xA5,0x00), 9.8,  "CRYPTO\nService"),
]
for color, lx, label in services:
    flow_box(s, lx, 2.85, 2.7, 0.7, label, CARD_BG, color, size=11)

# BATCH also connects to CRYPTO, STORAGE, BLOCKCHAIN, AI (design.md)
# Show as sub-row
sub_services = [
    (RGBColor(0xFF,0x69,0xB4), 3.6,  "AI Sentinel\nService"),
    (RGBColor(0x00,0xCE,0xD1), 6.7,  "STORAGE\nService"),
    (RGBColor(0xFF,0xFF,0x00), 9.8,  "BLOCKCHAIN\nService"),
]
for color, lx, label in sub_services:
    flow_box(s, lx, 3.75, 2.7, 0.65, label, CARD_BG, color, size=11)

# arrows: BATCH → AI, STORAGE, BLOCKCHAIN
for lx in [4.95, 8.05, 11.15]:
    arrow_down(s, lx, 3.58)

# arrows down to Data Layer
for lx in [1.25, 4.95, 8.05, 11.15]:
    arrow_down(s, lx, 4.43)

# ── Data Layer ────────────────────────────────────────────────────────────────
# PostgreSQL, AWS S3, Blockchain Network  (design.md data layer)
data = [
    (ACCENT,                   0.5,  "🗄️ PostgreSQL\nDatabase"),
    (ACCENT2,                  4.0,  "☁️ AWS S3\nBucket"),
    (GREEN,                    7.5,  "⛓️ Blockchain\nNetwork"),
    (RGBColor(0xFF,0xA5,0x00), 10.5, "📨 RabbitMQ\nJob Queue"),
]
for color, lx, label in data:
    flow_box(s, lx, 4.72, 2.6, 0.7, label, color, DARK_BG, size=11)

# bottom note
txt(s, "AUTH→PG  |  BATCH→PG  |  QR→PG  |  STORAGE→S3  |  BLOCKCHAIN→BC  |  AI→PG",
    0.5, 5.6, 12.33, 0.38, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Flowchart: Digital Signature Flow  (matches design.md §4.2)
# ══════════════════════════════════════════════════════════════════════════════
# Sequence: Mfg → API → CryptoSvc → DB → Blockchain
# Steps per design.md:
#   Mfg→API: Create Batch + Lab Report
#   API→CryptoSvc: generateHash(labReport) → SHA-256 Hash
#   API→CryptoSvc: signData(batchData, privateKey) → Digital Signature
#   API→DB: saveBatch(data, signature, hash)
#   API→BC: mintToken(batchData, hash, signature) → Transaction ID
#   API→DB: updateBatch(txId)
#   API→Mfg: Batch Created
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Digital Signature Flow", 0.5, 0.3, 12, 0.7,
    size=32, bold=True, color=WHITE)

# ── Participant headers (5 participants from design.md) ───────────────────────
participants = [
    (ACCENT,                   0.3,  "🏭\nManufacturer"),
    (ACCENT2,                  2.9,  "⚡\nAPI"),
    (GREEN,                    5.5,  "🔐\nCrypto\nService"),
    (RGBColor(0xFF,0xA5,0x00), 8.1,  "🗄️\nDB"),
    (RGBColor(0x00,0xCE,0xD1), 10.7, "⛓️\nBlockchain"),
]
for color, lx, label in participants:
    flow_box(s, lx, 1.0, 2.2, 0.75, label, color, DARK_BG, size=11)

# Lifelines (vertical dashed-style thin boxes)
lifeline_xs = [1.4, 4.0, 6.6, 9.2, 11.8]
for lx in lifeline_xs:
    box(s, lx - 0.01, 1.75, 0.02, 5.3, LIGHT_GRAY)

# ── Step 1: Mfg → API: Create Batch + Lab Report ─────────────────────────────
ty = 1.9
arrow_right(s, 1.42, ty + 0.05)
arrow_right(s, 1.9,  ty + 0.05)
arrow_right(s, 2.4,  ty + 0.05)
box(s, 1.42, ty, 2.55, 0.02, ACCENT)
txt(s, "Create Batch + Lab Report", 1.5, ty - 0.28, 2.4, 0.3,
    size=10, color=ACCENT)
txt(s, "▶", 3.85, ty - 0.1, 0.3, 0.3, size=12, color=ACCENT, align=PP_ALIGN.CENTER)

# ── Step 2: API → CryptoSvc: generateHash(labReport) ─────────────────────────
ty = 2.35
box(s, 4.02, ty, 2.55, 0.02, ACCENT2)
txt(s, "generateHash(labReport)", 4.1, ty - 0.28, 2.4, 0.3,
    size=10, color=ACCENT2)
txt(s, "▶", 6.45, ty - 0.1, 0.3, 0.3, size=12, color=ACCENT2, align=PP_ALIGN.CENTER)

# Return: SHA-256 Hash
ty = 2.72
box(s, 4.02, ty, 2.55, 0.02, GREEN)
txt(s, "SHA-256 Hash", 4.1, ty + 0.04, 2.4, 0.28,
    size=10, color=GREEN)
txt(s, "◀", 3.85, ty - 0.05, 0.3, 0.3, size=12, color=GREEN, align=PP_ALIGN.CENTER)

# ── Step 3: API → CryptoSvc: signData(batchData, privateKey) ─────────────────
ty = 3.1
box(s, 4.02, ty, 2.55, 0.02, ACCENT2)
txt(s, "signData(batchData, privateKey)", 4.1, ty - 0.28, 2.4, 0.3,
    size=10, color=ACCENT2)
txt(s, "▶", 6.45, ty - 0.1, 0.3, 0.3, size=12, color=ACCENT2, align=PP_ALIGN.CENTER)

# Return: Digital Signature
ty = 3.47
box(s, 4.02, ty, 2.55, 0.02, GREEN)
txt(s, "Digital Signature", 4.1, ty + 0.04, 2.4, 0.28,
    size=10, color=GREEN)
txt(s, "◀", 3.85, ty - 0.05, 0.3, 0.3, size=12, color=GREEN, align=PP_ALIGN.CENTER)

# ── Step 4: API → DB: saveBatch(data, signature, hash) ───────────────────────
ty = 3.85
box(s, 4.02, ty, 5.15, 0.02, RGBColor(0xFF,0xA5,0x00))
txt(s, "saveBatch(data, signature, hash)", 4.1, ty - 0.28, 5.0, 0.3,
    size=10, color=RGBColor(0xFF,0xA5,0x00))
txt(s, "▶", 9.05, ty - 0.1, 0.3, 0.3, size=12, color=RGBColor(0xFF,0xA5,0x00), align=PP_ALIGN.CENTER)

# ── Step 5: API → Blockchain: mintToken(batchData, hash, signature) ──────────
ty = 4.22
box(s, 4.02, ty, 7.75, 0.02, RGBColor(0x00,0xCE,0xD1))
txt(s, "mintToken(batchData, hash, signature)", 4.1, ty - 0.28, 7.6, 0.3,
    size=10, color=RGBColor(0x00,0xCE,0xD1))
txt(s, "▶", 11.65, ty - 0.1, 0.3, 0.3, size=12, color=RGBColor(0x00,0xCE,0xD1), align=PP_ALIGN.CENTER)

# Return: Transaction ID
ty = 4.6
box(s, 4.02, ty, 7.75, 0.02, GREEN)
txt(s, "Transaction ID (txId)", 7.5, ty + 0.04, 4.3, 0.28,
    size=10, color=GREEN)
txt(s, "◀", 3.85, ty - 0.05, 0.3, 0.3, size=12, color=GREEN, align=PP_ALIGN.CENTER)

# ── Step 6: API → DB: updateBatch(txId) ──────────────────────────────────────
ty = 4.97
box(s, 4.02, ty, 5.15, 0.02, RGBColor(0xFF,0xA5,0x00))
txt(s, "updateBatch(txId)", 4.1, ty - 0.28, 5.0, 0.3,
    size=10, color=RGBColor(0xFF,0xA5,0x00))
txt(s, "▶", 9.05, ty - 0.1, 0.3, 0.3, size=12, color=RGBColor(0xFF,0xA5,0x00), align=PP_ALIGN.CENTER)

# ── Step 7: API → Mfg: Batch Created ─────────────────────────────────────────
ty = 5.35
box(s, 1.42, ty, 2.55, 0.02, GREEN)
txt(s, "✅  Batch Created", 1.5, ty + 0.04, 2.4, 0.28,
    size=11, bold=True, color=GREEN)
txt(s, "◀", 1.25, ty - 0.05, 0.3, 0.3, size=12, color=GREEN, align=PP_ALIGN.CENTER)

# Bottom note
box(s, 0.4, 6.05, 12.53, 0.5, CARD_BG)
txt(s, "🔑 Private Key (Manufacturer) signs data  |  ⛓️ Blockchain stores immutable token  |  🗄️ DB stores signature + hash + txId",
    0.5, 6.13, 12.3, 0.35, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 8)


s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Authentication Flow", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

# Left column: Client side
txt(s, "CLIENT", 0.6, 1.15, 2.5, 0.3, size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
flow_box(s, 0.5, 1.45, 2.7, 0.5, "POST /auth/login\n(username + password)", CARD_BG, ACCENT)
arrow_down(s, 1.6, 2.0)
flow_box(s, 0.5, 2.3, 2.7, 0.5, "Receive JWT Token", CARD_BG, GREEN)
arrow_down(s, 1.6, 2.85)
flow_box(s, 0.5, 3.15, 2.7, 0.5, "GET /batches\n(Bearer Token)", CARD_BG, ACCENT)
arrow_down(s, 1.6, 3.7)
flow_box(s, 0.5, 4.0, 2.7, 0.5, "Receive Batch Data ✅", CARD_BG, GREEN)

# Arrows between columns
arrow_right(s, 3.25, 1.62)
arrow_right(s, 3.25, 3.32)

# Middle column: API Gateway
txt(s, "API GATEWAY", 3.7, 1.15, 2.8, 0.3, size=11, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
flow_box(s, 3.6, 1.45, 2.8, 0.5, "JWT Auth Filter\nValidate Token", CARD_BG, ACCENT2)
arrow_down(s, 4.9, 2.0)
flow_box(s, 3.6, 2.3, 2.8, 0.5, "Check Permission\n(RBAC)", CARD_BG, ACCENT2)
arrow_down(s, 4.9, 2.85)
flow_box(s, 3.6, 3.15, 2.8, 0.5, "Route to Service", CARD_BG, ACCENT2)
arrow_down(s, 4.9, 3.7)
flow_box(s, 3.6, 4.0, 2.8, 0.5, "Return Response", CARD_BG, ACCENT2)

arrow_right(s, 6.45, 1.62)
arrow_right(s, 6.45, 2.47)

# Right column: Auth Service + DB
txt(s, "AUTH SERVICE + DB", 6.9, 1.15, 3.5, 0.3, size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
flow_box(s, 6.8, 1.45, 3.5, 0.5, "findUser(username)\nfrom PostgreSQL", CARD_BG, GREEN)
arrow_down(s, 8.4, 2.0)
flow_box(s, 6.8, 2.3, 3.5, 0.5, "verifyPassword\n(BCrypt hash)", CARD_BG, GREEN)
arrow_down(s, 8.4, 2.85)
flow_box(s, 6.8, 3.15, 3.5, 0.5, "generateJWT\n(user + role claims)", CARD_BG, GREEN)
arrow_down(s, 8.4, 3.7)
flow_box(s, 6.8, 4.0, 3.5, 0.5, "Token → API → Client", CARD_BG, GREEN)

# Decision diamond for invalid
flow_box(s, 3.6, 5.0, 2.8, 0.6, "❌ Invalid Token?\n→ 401 Unauthorized", RED, WHITE, size=11)
txt(s, "If token invalid ↑", 3.6, 4.65, 2.8, 0.35, size=10, color=RED, align=PP_ALIGN.CENTER)

slide_number(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Flowchart: QR Code Security Flow
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "QR Code Security Flow", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

# Generation side (left)
txt(s, "QR GENERATION", 0.5, 1.1, 5.5, 0.35,
    size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

gen_steps = [
    (ACCENT,  "Batch Data + Serial Number"),
    (ACCENT,  "ECDSA Sign (Private Key)"),
    (ACCENT,  "Protobuf Encode → ~92 bytes"),
    (ACCENT,  "Base64 URL Encode"),
    (ACCENT,  "Generate QR Code Image"),
    (GREEN,   "Store in DB  |  max_scan=5  |  active=true"),
]
for i, (color, label) in enumerate(gen_steps):
    flow_box(s, 0.4, 1.5 + i * 0.82, 5.5, 0.65, label, CARD_BG, color, size=12)
    if i < len(gen_steps) - 1:
        arrow_down(s, 2.9, 2.18 + i * 0.82)

# Divider
box(s, 6.3, 1.1, 0.04, 6.0, ACCENT2)
txt(s, "VS", 6.1, 3.8, 0.5, 0.4, size=14, bold=True, color=ACCENT2)

# Verification side (right)
txt(s, "QR VERIFICATION", 6.8, 1.1, 6.0, 0.35,
    size=13, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

ver_steps = [
    (ACCENT2, "Scan QR → Extract Base64 URL"),
    (ACCENT2, "Decode Protobuf Payload"),
    (ACCENT2, "Verify ECDSA Signature"),
    (RGBColor(0xFF,0xA5,0x00), "Check: scan_count < max_scan_limit"),
    (RGBColor(0xFF,0xA5,0x00), "Check: Geographic Anomaly (AI)"),
]
for i, (color, label) in enumerate(ver_steps):
    flow_box(s, 6.8, 1.5 + i * 0.82, 6.0, 0.65, label, CARD_BG, color, size=12)
    if i < len(ver_steps) - 1:
        arrow_down(s, 9.7, 2.18 + i * 0.82)

# Result boxes
flow_box(s, 6.8, 5.65, 2.7, 0.65, "✅  VALID\nAuthentic Medicine", GREEN, DARK_BG, size=12)
flow_box(s, 10.1, 5.65, 2.7, 0.65, "❌  FRAUD\nAuto Kill-Switch", RED, WHITE, size=12)
txt(s, "Pass ↙          Fail ↘", 6.8, 5.35, 6.0, 0.3,
    size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Flowchart: Supply Chain Flow
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Supply Chain Flow", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

chain = [
    (ACCENT,  "🏭", "MANUFACTURER",   "Create Batch\nUpload Lab Report\nSign + Blockchain"),
    (ACCENT2, "✅", "QC APPROVAL",    "Multi-Sig Approval\nProduction Head\n+ Quality Checker"),
    (GREEN,   "📦", "DISTRIBUTOR",    "Accept Transfer\nScan & Verify\nUpdate Ownership"),
    (RGBColor(0xFF,0xA5,0x00), "🏪", "RETAILER", "Receive Stock\nVerify QR\nTrack Inventory"),
    (RGBColor(0xFF,0x69,0xB4), "💊", "PATIENT",  "Scan QR\nVerify Authenticity\nCheck Expiry"),
]

for i, (color, icon, role, desc) in enumerate(chain):
    lx = 0.3 + i * 2.55
    # role box
    flow_box(s, lx, 1.3, 2.3, 0.55, f"{icon}  {role}", color, DARK_BG, size=12)
    # desc box
    box(s, lx, 1.9, 2.3, 1.6, CARD_BG)
    box(s, lx, 1.9, 0.04, 1.6, color)
    for j, line in enumerate(desc.split('\n')):
        txt(s, line, lx + 0.12, 1.98 + j * 0.48, 2.1, 0.45, size=11, color=LIGHT_GRAY)
    # arrow
    if i < len(chain) - 1:
        txt(s, "▶", lx + 2.35, 1.48, 0.4, 0.35, size=16, color=ACCENT, align=PP_ALIGN.CENTER)

# Blockchain layer below
box(s, 0.3, 3.75, 12.7, 0.05, ACCENT2)
txt(s, "⛓️  BLOCKCHAIN LAYER  —  Immutable Record for Every Step",
    0.3, 3.85, 12.7, 0.45, size=13, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# AI Sentinel layer
box(s, 0.3, 4.45, 12.7, 1.3, CARD_BG)
box(s, 0.3, 4.45, 12.7, 0.05, RED)
txt(s, "🤖  AI SENTINEL  —  Monitoring Every Scan in Real-Time",
    0.5, 4.55, 12.3, 0.38, size=13, bold=True, color=RED)
sentinel_items = [
    "Impossible Travel Detection  →  Auto Kill-Switch",
    "Simultaneous Multi-Location Scans  →  Immediate Recall",
    "Device Fingerprint Mismatch  →  Alert + Manual Review",
    "Scan Count Exceeded  →  Unit Deactivated",
]
for i, item in enumerate(sentinel_items):
    col = i % 2
    row = i // 2
    txt(s, f"⚡ {item}", 0.5 + col * 6.5, 5.0 + row * 0.38, 6.3, 0.35,
        size=11, color=LIGHT_GRAY)

slide_number(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Key Features", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

features = [
    (ACCENT,  "Multi-Sig Batch Approval",
     ["Production Head + Quality Checker", "Digital signatures on every batch", "Blockchain token minted on approval"]),
    (ACCENT2, "AI Anomaly Detection",
     ["Impossible travel (>200 km/h)", "Simultaneous multi-location scans", "Auto kill-switch triggers instantly"]),
    (GREEN,   "QR Anti-Cloning",
     ["ECDSA signed compact payload", "Max scan limit enforcement", "Offline signature verification"]),
    (RGBColor(0xFF,0xA5,0x00), "Granular RBAC",
     ["Regulator sees only alerts/recalls", "Manufacturer data stays private", "Full audit trail for compliance"]),
    (RGBColor(0xFF,0x69,0xB4), "Async Job Processing",
     ["RabbitMQ background workers", "Idempotent crash-safe generation", "Real-time progress tracking"]),
    (RGBColor(0x00,0xCE,0xD1), "Expiry Monitoring",
     ["Automated expiry alerts", "Scheduled batch scanning", "Proactive recall initiation"]),
]
for i, (color, title, bullets) in enumerate(features):
    col = i % 3
    row = i // 3
    lx = 0.4 + col * 4.3
    ty = 1.3 + row * 2.85
    box(s, lx, ty, 4.0, 2.6, CARD_BG)
    box(s, lx, ty, 0.05, 2.6, color)
    txt(s, title, lx + 0.15, ty + 0.1, 3.7, 0.4,
        size=13, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", lx + 0.15, ty + 0.6 + j * 0.55, 3.7, 0.5,
            size=11, color=LIGHT_GRAY)

slide_number(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Tech Stack", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

stack = [
    ("Backend",    ACCENT,  ["Java 17 + Spring Boot 3", "Spring Security + JWT", "JPA / Hibernate"]),
    ("Database",   ACCENT2, ["PostgreSQL (production)", "Flyway migrations", "Redis caching"]),
    ("Frontend",   GREEN,   ["React + Vite", "Tailwind CSS", "Axios REST client"]),
    ("Cloud",      RGBColor(0xFF,0xA5,0x00), ["AWS S3 (lab reports)", "Docker + Compose", "RabbitMQ queues"]),
    ("Security",   RGBColor(0xFF,0x69,0xB4), ["ECDSA digital signatures", "SHA-256 hashing", "Protobuf QR payloads"]),
    ("Blockchain", RGBColor(0x00,0xCE,0xD1), ["Ethereum smart contracts", "Merkle root storage", "Web3j integration"]),
]
for i, (cat, color, items) in enumerate(stack):
    col = i % 3
    row = i // 2
    lx = 0.4 + col * 4.3
    ty = 1.3 + row * 2.7
    box(s, lx, ty, 4.0, 2.4, CARD_BG)
    box(s, lx, ty, 4.0, 0.45, color)
    txt(s, cat, lx + 0.15, ty + 0.05, 3.7, 0.38,
        size=15, bold=True, color=DARK_BG)
    for j, item in enumerate(items):
        txt(s, f"  ▸  {item}", lx + 0.1, ty + 0.6 + j * 0.55, 3.7, 0.5,
            size=12, color=WHITE)

slide_number(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Security Deep Dive
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Security Architecture", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

# Left: Auth flow
box(s, 0.4, 1.2, 5.9, 5.8, CARD_BG)
box(s, 0.4, 1.2, 0.05, 5.8, ACCENT)
txt(s, "Authentication Flow", 0.6, 1.3, 5.5, 0.4,
    size=14, bold=True, color=ACCENT)
auth_steps = [
    "1. Client → POST /auth/login",
    "2. AuthService verifies BCrypt hash",
    "3. JWT token generated with role claims",
    "4. All requests validated via JWT filter",
    "5. RBAC checks on every endpoint",
    "6. Token blacklist on logout",
]
for i, step in enumerate(auth_steps):
    txt(s, step, 0.6, 1.85 + i * 0.72, 5.5, 0.6,
        size=12, color=LIGHT_GRAY)

# Right: QR security
box(s, 6.9, 1.2, 5.9, 5.8, CARD_BG)
box(s, 6.9, 1.2, 0.05, 5.8, ACCENT2)
txt(s, "QR Anti-Cloning Security", 7.1, 1.3, 5.5, 0.4,
    size=14, bold=True, color=ACCENT2)
qr_steps = [
    "1. ECDSA sign: serial + batch + timestamp",
    "2. Protobuf encode → ~92 bytes payload",
    "3. Base64 URL → QR code generated",
    "4. Max 5 scans per unit enforced",
    "5. Geographic anomaly → auto kill-switch",
    "6. Offline verify via cached public key",
]
for i, step in enumerate(qr_steps):
    txt(s, step, 7.1, 1.85 + i * 0.72, 5.5, 0.6,
        size=12, color=LIGHT_GRAY)

slide_number(s, 12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Blockchain Integration
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Blockchain Integration", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

# Key insight box
box(s, 0.5, 1.2, 12.33, 1.1, CARD_BG)
box(s, 0.5, 1.2, 12.33, 0.05, ACCENT)
txt(s, "💡  Key Innovation: Only Merkle Root stored on-chain → 99.99% gas cost reduction",
    0.7, 1.35, 11.9, 0.6, size=16, bold=True, color=ACCENT)

points = [
    (ACCENT,  "What's On-Chain",
     "Batch number · Medicine hash · Expiry date · Manufacturer address · Merkle root · Total units"),
    (ACCENT2, "What's Off-Chain",
     "Individual unit data · Lab reports (S3) · Scan logs · Business metrics"),
    (GREEN,   "Recall Events",
     "Emitted as blockchain events · Immutable audit trail · Auto-triggered by AI Sentinel"),
    (RGBColor(0xFF,0xA5,0x00), "Verification Flow",
     "QR scan → Merkle proof → Blockchain root match → Cryptographic guarantee of authenticity"),
]
for i, (color, title, body) in enumerate(points):
    ty = 2.55 + i * 1.15
    box(s, 0.5, ty, 12.33, 1.0, CARD_BG)
    box(s, 0.5, ty, 0.05, 1.0, color)
    txt(s, title, 0.7, ty + 0.08, 3.0, 0.38,
        size=13, bold=True, color=color)
    txt(s, body, 3.9, ty + 0.2, 8.7, 0.6,
        size=12, color=LIGHT_GRAY)

slide_number(s, 13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — User Roles & Dashboard
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "User Roles & Dashboard", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

roles = [
    (ACCENT,  "🏭 Manufacturer",
     ["Create & approve batches", "Upload lab reports", "Generate QR codes", "Monitor expiry alerts"]),
    (ACCENT2, "🚚 Distributor",
     ["Accept unit transfers", "Scan & verify QR codes", "Track ownership chain", "Report anomalies"]),
    (GREEN,   "💊 Pharmacist / Patient",
     ["Scan QR to verify medicine", "View batch details", "Check expiry & recalls", "Offline verification"]),
    (RGBColor(0xFF,0xA5,0x00), "🏛️ Regulator",
     ["View fraud alerts only", "Initiate emergency recalls", "Access audit logs", "No business metrics"]),
]
for i, (color, role, bullets) in enumerate(roles):
    lx = 0.4 + i * 3.15
    box(s, lx, 1.3, 2.9, 5.6, CARD_BG)
    box(s, lx, 1.3, 2.9, 0.06, color)
    txt(s, role, lx + 0.1, 1.45, 2.7, 0.5,
        size=14, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", lx + 0.1, 2.1 + j * 0.65, 2.7, 0.6,
            size=12, color=LIGHT_GRAY)

slide_number(s, 14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Impact & Metrics
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Impact & Key Metrics", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

metrics = [
    (GREEN,   "99.99%",  "Blockchain gas cost reduction\nvia Merkle root optimization"),
    (ACCENT,  "71%",     "QR payload size reduction\nProtobuf vs JSON (320→92 bytes)"),
    (ACCENT2, "< 5 sec", "Fraud detection response time\nAI Sentinel auto kill-switch"),
    (RGBColor(0xFF,0xA5,0x00), "5 Roles", "Complete RBAC coverage\nManufacturer to Regulator"),
    (RGBColor(0xFF,0x69,0xB4), "100%",    "Crash-safe job processing\nIdempotent unit generation"),
    (RGBColor(0x00,0xCE,0xD1), "0",       "Single point of failure\nDistributed cloud architecture"),
]
for i, (color, metric, desc) in enumerate(metrics):
    col = i % 3
    row = i // 3
    lx = 0.4 + col * 4.3
    ty = 1.3 + row * 2.8
    box(s, lx, ty, 4.0, 2.5, CARD_BG)
    box(s, lx, ty, 4.0, 0.06, color)
    txt(s, metric, lx + 0.15, ty + 0.2, 3.7, 0.9,
        size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc, lx + 0.1, ty + 1.15, 3.8, 1.1,
        size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Live Demo Flow
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Demo Flow", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

steps = [
    (ACCENT,  "1", "Manufacturer Login",    "JWT auth → role-based dashboard loads"),
    (ACCENT2, "2", "Create Batch",          "Upload lab report → S3 storage → SHA-256 hash"),
    (GREEN,   "3", "Multi-Sig Approval",    "Production Head + QC sign → blockchain token minted"),
    (RGBColor(0xFF,0xA5,0x00), "4", "QR Generation", "Async job → 10,000 units → Protobuf QR codes"),
    (RGBColor(0xFF,0x69,0xB4), "5", "Scan & Verify",  "Pharmacist scans → AI checks → VALID / FRAUD"),
    (RGBColor(0x00,0xCE,0xD1), "6", "Fraud Simulation","Impossible travel → auto kill-switch → regulator alert"),
]
for i, (color, num, title, desc) in enumerate(steps):
    ty = 1.3 + i * 0.98
    box(s, 0.4, ty, 12.33, 0.85, CARD_BG)
    # number circle
    box(s, 0.4, ty, 0.85, 0.85, color)
    txt(s, num, 0.4, ty + 0.15, 0.85, 0.55,
        size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, title, 1.4, ty + 0.08, 3.5, 0.38,
        size=14, bold=True, color=color)
    txt(s, desc, 5.1, ty + 0.2, 7.4, 0.45,
        size=13, color=LIGHT_GRAY)

slide_number(s, 16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Future Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

txt(s, "Future Roadmap", 0.5, 0.3, 12, 0.7,
    size=36, bold=True, color=WHITE)

# Government Integration highlight box
box(s, 0.4, 1.15, 12.53, 1.05, CARD_BG)
box(s, 0.4, 1.15, 12.53, 0.05, GREEN)
txt(s, "🏛️  Government Integration (Priority)",
    0.65, 1.22, 5.5, 0.38, size=14, bold=True, color=GREEN)
txt(s, "Planning to integrate with ABHA (Ayushman Bharat Health Account) — linking medicine verification",
    0.65, 1.6, 12.0, 0.38, size=12, color=LIGHT_GRAY)
txt(s, "directly to patient health records for seamless, government-backed pharmaceutical traceability.",
    0.65, 1.92, 12.0, 0.35, size=12, color=LIGHT_GRAY)

# 6 future items in 2 columns
future_items = [
    (ACCENT,  "🏥  ABHA / NHA Integration",
     ["Link medicine scans to ABHA health ID", "NHA API integration for patient records", "PM-JAY scheme medicine verification"]),
    (ACCENT2, "🤖  Advanced AI / ML",
     ["Predictive counterfeit detection model", "NLP-based complaint analysis", "Computer vision for label verification"]),
    (GREEN,   "📱  Mobile App (Android/iOS)",
     ["Native QR scanner for patients", "Offline-first architecture", "Multi-language support (Hindi, regional)"]),
    (RGBColor(0xFF,0xA5,0x00), "🌐  National Drug Registry",
     ["Connect with CDSCO drug database", "Auto-validate approved medicines", "Real-time recall notifications"]),
    (RGBColor(0xFF,0x69,0xB4), "⛓️  Multi-Chain Blockchain",
     ["Hyperledger for enterprise partners", "Cross-chain batch verification", "Smart contract auto-compliance"]),
    (RGBColor(0x00,0xCE,0xD1), "📊  Analytics Dashboard",
     ["National supply chain heatmap", "Counterfeit hotspot detection", "Regulator BI reports & exports"]),
]
for i, (color, title, bullets) in enumerate(future_items):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.5
    ty = 2.45 + row * 1.65
    box(s, lx, ty, 6.1, 1.5, CARD_BG)
    box(s, lx, ty, 0.05, 1.5, color)
    txt(s, title, lx + 0.15, ty + 0.08, 5.7, 0.38,
        size=13, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", lx + 0.15, ty + 0.52 + j * 0.32, 5.7, 0.3,
            size=11, color=LIGHT_GRAY)

slide_number(s, 17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0, 0, 13.33, 0.18, ACCENT)
box(s, 6.67, 0, 6.66, 0.18, ACCENT2)

txt(s, "Thank You", 1, 1.5, 11, 1.2,
    size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 3.5, 3.1, 6.33, 0.04, ACCENT2)

txt(s, "PharmaTrust AI  —  Securing Every Pill, Every Step",
    1, 3.3, 11, 0.6, size=20, color=ACCENT, align=PP_ALIGN.CENTER)

contact_items = [
    "🔗  Blockchain-verified supply chain",
    "🤖  AI-powered fraud detection",
    "📱  QR anti-cloning technology",
]
for i, item in enumerate(contact_items):
    txt(s, item, 3.5, 4.2 + i * 0.55, 6.33, 0.5,
        size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "Questions?", 1, 6.5, 11, 0.5,
    size=18, color=ACCENT2, align=PP_ALIGN.CENTER)
slide_number(s, 18)


# ── Save ───────────────────────────────────────────────────────────────────────
output = "PharmaTrust_Hackathon.pptx"
prs.save(output)
print(f"✅  Presentation saved: {output}  ({prs.slides.__len__()} slides)")
